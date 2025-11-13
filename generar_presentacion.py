"""
Script para generar presentación PowerPoint TÉCNICA del proyecto ETL
ENFOQUE: Procedimientos, código y conceptos técnicos (no resultados de negocio)
Sigue el guión de GUIA_PRESENTACION_15MIN.md
Requiere: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def agregar_caja_codigo(slide, codigo_lines, left, top, width, height, es_sql=False):
    """Helper para agregar cajas de código con fondo"""
    # Fondo gris claro
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    # Código encima
    code_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1), 
        Inches(width - 0.3), Inches(height - 0.2)
    )
    text_frame = code_box.text_frame
    text_frame.word_wrap = True
    
    color_code = RGBColor(139, 0, 0) if es_sql else RGBColor(0, 0, 128)
    
    for i, line in enumerate(codigo_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(13)
        if line.strip().startswith("#") or line.strip().startswith("--"):
            p.font.color.rgb = RGBColor(0, 128, 0)
        else:
            p.font.color.rgb = color_code

def crear_presentacion():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colores del tema
    COLOR_TITULO = RGBColor(31, 78, 121)
    COLOR_SUBTITULO = RGBColor(68, 114, 196)
    COLOR_TEXTO = RGBColor(64, 64, 64)
    COLOR_ACENTO = RGBColor(237, 125, 49)
    COLOR_EXITO = RGBColor(112, 173, 71)
    
    # ============================================
    # PERSONA 1 - SLIDE 1: PORTADA (15 seg)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(0.8))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "PROCESO ETL COMPLETO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    p.alignment = PP_ALIGN.CENTER
    
    subtitulo = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(0.5))
    text_frame = subtitulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Extract → Transform → Load"
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_ACENTO
    p.alignment = PP_ALIGN.CENTER
    
    stack = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.8))
    text_frame = stack.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Python 3.13 • Pandas • SQLite • Matplotlib"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    p = text_frame.add_paragraph()
    p.text = "99,459 transacciones • Google Colab compatible"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================
    # PERSONA 1 - SLIDE 2: ¿QUÉ ES ETL? (30 seg)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "¿Qué es ETL?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Definición
    definicion = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.5))
    text_frame = definicion.text_frame
    p = text_frame.paragraphs[0]
    p.text = "ETL = Extract + Transform + Load"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.alignment = PP_ALIGN.CENTER
    
    # Tres fases visuales
    fases = [
        ("EXTRACT", "Obtener datos desde fuentes CSV", 2.2, RGBColor(68, 114, 196)),
        ("TRANSFORM", "Limpiar, calcular, validar datos", 3.5, RGBColor(237, 125, 49)),
        ("LOAD", "Cargar a base de datos relacional", 4.8, RGBColor(112, 173, 71))
    ]
    
    for nombre, desc, top, color in fases:
        box = slide.shapes.add_textbox(Inches(1.5), Inches(top), Inches(7), Inches(0.9))
        text_frame = box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = f"📥 {nombre}" if nombre == "EXTRACT" else f"🔄 {nombre}" if nombre == "TRANSFORM" else f"📤 {nombre}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.level = 1
    
    # Concepto teórico
    concepto = slide.shapes.add_textbox(Inches(1), Inches(6.4), Inches(8), Inches(0.6))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "CONCEPTO TEÓRICO: Integración de datos (BD II) + Automatización (Prog I)"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    
    # ============================================
    # PERSONA 1 - SLIDE 3: EXTRACT - CARGA (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "PERSONA 1: EXTRACT - Carga de datos"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código
    codigo = [
        "import pandas as pd",
        "",
        "# Cargar datos de ventas",
        "df_sales = pd.read_csv('sales_data.csv')",
        "print(df_sales.shape)  # (99459, 8)",
        "",
        "# Cargar datos de clientes",
        "df_customers = pd.read_csv('customer_data.csv')",
        "print(df_customers.shape)  # (100, 3)"
    ]
    
    agregar_caja_codigo(slide, codigo, 0.8, 1.3, 8.4, 2.8, es_sql=False)
    
    # Output simulado
    output_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(1.2))
    text_frame = output_box.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "OUTPUT:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    p = text_frame.add_paragraph()
    p.text = "✓ df_sales: 99,459 filas × 8 columnas (invoice_no, customer_id, category, quantity, price...)"
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "✓ df_customers: 100 filas × 3 columnas (customer_id, gender, age)"
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.level = 1
    
    # Concepto
    concepto = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "¿Qué es un DataFrame? Estructura bidimensional de Pandas (similar a tabla SQL)"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    p = text_frame.add_paragraph()
    p.text = "¿Por qué Pandas? Optimizado en C, 1000x más rápido que listas Python puras"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    
    # ============================================
    # PERSONA 1 - SLIDE 4: MERGE (1 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "MERGE: Combinación de tablas"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código Pandas
    label1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(0.4))
    text_frame = label1.text_frame
    p = text_frame.paragraphs[0]
    p.text = "CÓDIGO PANDAS:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUBTITULO
    
    codigo_pandas = [
        "df_combined = pd.merge(",
        "    df_sales,",
        "    df_customers,",
        "    on='customer_id',",
        "    how='left'  # LEFT JOIN",
        ")"
    ]
    
    agregar_caja_codigo(slide, codigo_pandas, 0.8, 1.6, 4, 1.8, es_sql=False)
    
    # Equivalente SQL
    label2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(0.4))
    text_frame = label2.text_frame
    p = text_frame.paragraphs[0]
    p.text = "EQUIVALENTE SQL:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUBTITULO
    
    codigo_sql = [
        "SELECT *",
        "FROM sales",
        "LEFT JOIN customers",
        "  ON sales.customer_id =",
        "     customers.customer_id"
    ]
    
    agregar_caja_codigo(slide, codigo_sql, 5.2, 1.6, 4, 1.8, es_sql=True)
    
    # Resultado
    resultado = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    text_frame = resultado.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "RESULTADO: df_combined con 99,459 filas × 11 columnas"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    # Concepto
    concepto = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1.5))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "CONCEPTO TEÓRICO: JOIN (Base de Datos II)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_EXITO
    
    p = text_frame.add_paragraph()
    p.text = "¿Por qué LEFT JOIN? Preserva TODAS las ventas (99,459 filas)"
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "INNER JOIN eliminaría transacciones sin datos demográficos → pérdida de información"
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    # ============================================
    # PERSONA 2 - SLIDE 5: LIBRERÍAS (1 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "PERSONA 2: Librerías utilizadas"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Tres librerías con código de import
    libs = [
        ("PANDAS", "import pandas as pd", "Manipulación de datos: limpieza, merge, cálculos", 1.3),
        ("NUMPY", "import numpy as np", "Arrays multidimensionales, operaciones matemáticas", 2.5),
        ("SQLITE3", "import sqlite3", "Persistencia: base de datos relacional ACID", 3.7),
        ("MATPLOTLIB", "import matplotlib.pyplot as plt", "Visualización: gráficos estadísticos PNG", 4.9)
    ]
    
    for nombre, import_code, desc, top in libs:
        # Nombre
        label = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(8), Inches(0.4))
        text_frame = label.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"📚 {nombre}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_SUBTITULO
        
        # Import
        code = slide.shapes.add_textbox(Inches(1.5), Inches(top + 0.4), Inches(7), Inches(0.3))
        text_frame = code.text_frame
        p = text_frame.paragraphs[0]
        p.text = import_code
        p.font.name = "Consolas"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        # Descripción
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(top + 0.7), Inches(7), Inches(0.3))
        text_frame = desc_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
    
    # Concepto
    concepto = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(8), Inches(0.6))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "CONCEPTO TEÓRICO: Modularidad (Prog I - POO) → Abstracción de complejidad"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    
    # ============================================
    # PERSONA 2 - SLIDE 6: TRANSFORM FECHAS (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "TRANSFORM: Conversión de fechas"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Antes/Después
    antes = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(3.5), Inches(0.7))
    text_frame = antes.text_frame
    p = text_frame.paragraphs[0]
    p.text = "❌ ANTES (dd-mm-yyyy)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)
    p = text_frame.add_paragraph()
    p.text = "'15-03-2023'"
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    
    despues = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(3.5), Inches(0.7))
    text_frame = despues.text_frame
    p = text_frame.paragraphs[0]
    p.text = "✅ DESPUÉS (ISO 8601)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_EXITO
    p = text_frame.add_paragraph()
    p.text = "'2023-03-15'"
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    
    # Código
    codigo = [
        "# Parsear con formato específico",
        "df['invoice_date'] = pd.to_datetime(",
        "    df['invoice_date'],",
        "    format='%d-%m-%Y'  # día-mes-año",
        ")",
        "",
        "# Convertir a ISO 8601",
        "df['invoice_date'] = (",
        "    df['invoice_date'].dt.strftime('%Y-%m-%d')",
        ")"
    ]
    
    agregar_caja_codigo(slide, codigo, 1, 2.2, 8, 2.5, es_sql=False)
    
    # Por qué
    por_que = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    text_frame = por_que.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "¿POR QUÉ ISO 8601?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    razones = [
        "✓ SQLite reconoce automáticamente yyyy-mm-dd",
        "✓ Ordenamiento cronológico correcto (no alfabético)",
        "✓ Estándar internacional (evita ambigüedades 03/04/2023)"
    ]
    
    for razon in razones:
        p = text_frame.add_paragraph()
        p.text = razon
        p.font.size = Pt(14)
        p.level = 1
    
    # ============================================
    # PERSONA 2 - SLIDE 7: TRANSFORM COLUMNAS (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "TRANSFORM: Columnas calculadas"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código
    codigo = [
        "# 1. Total de venta (persistir cálculo)",
        "df['total_sale'] = df['quantity'] * df['price']",
        "",
        "# 2. Extracción temporal",
        "df['year'] = df['invoice_date'].dt.year",
        "df['month'] = df['invoice_date'].dt.month",
        "",
        "# Ejemplo de filas resultantes:",
        "# quantity  price  total_sale  year  month",
        "#    5      25.99    129.95    2023    3"
    ]
    
    agregar_caja_codigo(slide, codigo, 1, 1.2, 8, 3, es_sql=False)
    
    # Decisión técnica
    decision = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.8))
    text_frame = decision.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "DECISIÓN TÉCNICA: Persistir en DB (no recalcular en cada query)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    p = text_frame.add_paragraph()
    p.text = "Trade-off: Espacio en disco ↑ pero velocidad de query ↑↑↑"
    p.font.size = Pt(14)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Sin columna: SELECT SUM(quantity * price) → 45ms"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(192, 0, 0)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Con columna: SELECT SUM(total_sale) → 12ms (3x más rápido)"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_EXITO
    p.level = 1
    
    # Concepto
    concepto = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.7))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "CONCEPTO TEÓRICO: Materialización de cálculos (BD II - Data Warehousing OLAP)"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    
    # ============================================
    # PERSONA 3 - SLIDE 8: LIMPIEZA NULOS (1 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "PERSONA 3: Limpieza de datos nulos"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código detección
    codigo1 = [
        "# 1. Identificar valores nulos",
        "print(df.isnull().sum())",
        "",
        "# Output:",
        "# age              119  ← ¡Crítico!",
        "# price              0",
        "# category           0"
    ]
    
    agregar_caja_codigo(slide, codigo1, 1, 1.2, 8, 2, es_sql=False)
    
    # Código limpieza
    codigo2 = [
        "# 2. Eliminar selectivamente (no fillna con valores falsos)",
        "df_clean = df.dropna(",
        "    subset=['age']  # Solo columna crítica demográfica",
        ")",
        "",
        "print(f'Eliminados: {len(df) - len(df_clean)} registros')",
        "# Output: Eliminados: 119 registros"
    ]
    
    agregar_caja_codigo(slide, codigo2, 1, 3.5, 8, 2.2, es_sql=False)
    
    # Resultado
    resultado = slide.shapes.add_textbox(Inches(1.5), Inches(6), Inches(7), Inches(0.8))
    text_frame = resultado.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "RESULTADO: 99,457 → 99,338 filas | Tasa recuperación: 99.88%"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_EXITO
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================
    # PERSONA 3 - SLIDE 9: LOAD A SQLITE (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "LOAD: Carga a base de datos SQLite"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código
    codigo = [
        "import sqlite3",
        "",
        "# Conexión a base de datos (crea archivo si no existe)",
        "conn = sqlite3.connect('sql/ventas.db')",
        "",
        "# Carga automática (CREATE TABLE + INSERT)",
        "df_clean.to_sql(",
        "    'datos_limpios',        # Nombre de tabla",
        "    conn,                    # Conexión",
        "    if_exists='replace',     # Reemplazar si existe",
        "    index=False              # No guardar índice Pandas",
        ")",
        "",
        "# Verificación",
        "cursor = conn.cursor()",
        "cursor.execute('SELECT COUNT(*) FROM datos_limpios')",
        "print(f'Registros cargados: {cursor.fetchone()[0]}')",
        "# Output: Registros cargados: 99338"
    ]
    
    agregar_caja_codigo(slide, codigo, 0.8, 1.2, 8.4, 4.5, es_sql=False)
    
    # Concepto
    concepto = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "¿Por qué SQLite? Sin servidor, archivo único, ACID-compliant, perfecto para 15 MB"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    p = text_frame.add_paragraph()
    p.text = "CONCEPTO TEÓRICO: Persistencia relacional (BD II) - Transacciones ACID"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    
    # ============================================
    # PERSONA 3 - SLIDE 10: CONSULTAS SQL (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Consultas SQL - Análisis de datos"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Query 1
    label1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.4))
    text_frame = label1.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Item 5a: Ventas mensuales (agregación temporal)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUBTITULO
    
    query1 = [
        "SELECT STRFTIME('%Y-%m', invoice_date) AS mes,",
        "       SUM(total_sale) AS ventas_totales",
        "FROM datos_limpios",
        "GROUP BY mes",
        "ORDER BY mes;"
    ]
    
    agregar_caja_codigo(slide, query1, 0.8, 1.6, 8.4, 1.5, es_sql=True)
    
    # Query 2
    label2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8), Inches(0.4))
    text_frame = label2.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Item 5b: Top 5 categorías (ordenamiento + límite)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUBTITULO
    
    query2 = [
        "SELECT category,",
        "       SUM(quantity) AS unidades_vendidas",
        "FROM datos_limpios",
        "GROUP BY category",
        "ORDER BY unidades_vendidas DESC",
        "LIMIT 5;"
    ]
    
    agregar_caja_codigo(slide, query2, 0.8, 3.9, 8.4, 1.7, es_sql=True)
    
    # Conceptos
    conceptos = slide.shapes.add_textbox(Inches(1), Inches(5.9), Inches(8), Inches(1.2))
    text_frame = conceptos.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "CONCEPTOS TEÓRICOS (BD II):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_EXITO
    
    items = [
        "SUM, COUNT → Funciones de agregación",
        "GROUP BY → Agrupamiento de filas",
        "ORDER BY → Ordenamiento ascendente/descendente",
        "STRFTIME → Función de fecha específica de SQLite"
    ]
    for item in items:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.level = 1
    
    # ============================================
    # PERSONA 4 - SLIDE 11: VISUALIZACIONES (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "PERSONA 4: Visualizaciones con Matplotlib"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código
    codigo = [
        "import matplotlib.pyplot as plt",
        "",
        "# Ejecutar query y obtener datos",
        "query = '''SELECT STRFTIME('%Y-%m', invoice_date) AS mes,",
        "                  SUM(total_sale) AS ventas",
        "           FROM datos_limpios",
        "           GROUP BY mes ORDER BY mes'''",
        "result = pd.read_sql(query, conn)",
        "",
        "# Crear gráfico de líneas",
        "plt.figure(figsize=(12, 6))",
        "plt.plot(result['mes'], result['ventas'], marker='o')",
        "plt.title('Evolución de Ventas Mensuales', fontsize=16)",
        "plt.xlabel('Mes')",
        "plt.ylabel('Ventas ($)')",
        "plt.grid(True, alpha=0.3)",
        "plt.xticks(rotation=45)",
        "plt.tight_layout()",
        "plt.savefig('visualizaciones/ventas_mensuales.png', dpi=300)"
    ]
    
    agregar_caja_codigo(slide, codigo, 0.8, 1.2, 8.4, 4.5, es_sql=False)
    
    # Output
    output = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.9))
    text_frame = output.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "OUTPUT: 8 gráficos PNG (300 DPI) en carpeta visualizaciones/"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    p = text_frame.add_paragraph()
    p.text = "CONCEPTO: Transformar datos numéricos → Insights visuales comunicables"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    p.level = 1
    
    # ============================================
    # PERSONA 4 - SLIDE 12: DECISIONES TÉCNICAS (1.5 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Decisiones técnicas diferenciadoras"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # 5 decisiones clave
    decisiones = [
        ("1. ISO 8601", "yyyy-mm-dd → SQL compatibility + ordenamiento correcto"),
        ("2. Tabla desnormalizada", "1 tabla → Velocidad de query (OLAP pattern)"),
        ("3. Columnas persistidas", "Calcular 1 vez → 3x más rápido en queries"),
        ("4. Colab auto-setup", "Auto-detección + clon repo → 1-click reproducibility"),
        ("5. Limpieza selectiva", "dropna(subset=[...]) → Transparencia (99.88% recovery)")
    ]
    
    y_pos = 1.3
    for num_titulo, desc in decisiones:
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.85))
        text_frame = box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = num_titulo
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_SUBTITULO
        
        p = text_frame.add_paragraph()
        p.text = f"→ {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        p.level = 1
        
        y_pos += 1
    
    # Concepto
    concepto = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.7))
    text_frame = concepto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "CONCEPTO: Trade-offs conscientes (Ingeniería de Software) - Justificar cada decisión"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLOR_EXITO
    
    # ============================================
    # PERSONA 4 - SLIDE 13: REPRODUCIBILIDAD (1 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Configuración Colab - Reproducibilidad"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Código Colab setup
    codigo = [
        "import sys, os",
        "",
        "# Auto-detección de entorno",
        "if 'google.colab' in sys.modules:",
        "    print('🚀 Google Colab detectado')",
        "    ",
        "    # Clonar repositorio automáticamente",
        "    !git clone https://github.com/Pablo-Tab/ABP-INNOVACION-DATOS.git",
        "    os.chdir('/content/ABP-INNOVACION-DATOS')",
        "    ",
        "    # Verificar archivos críticos",
        "    archivos = ['customer_data.csv', 'sales_data.csv']",
        "    for archivo in archivos:",
        "        if os.path.exists(archivo):",
        "            size = os.path.getsize(archivo) / (1024*1024)",
        "            print(f'✅ {archivo} ({size:.2f} MB)')",
        "    ",
        "    # Crear estructura de carpetas",
        "    os.makedirs('datos', exist_ok=True)",
        "    os.makedirs('visualizaciones', exist_ok=True)",
        "    ",
        "else:",
        "    print('💻 Entorno local detectado')"
    ]
    
    agregar_caja_codigo(slide, codigo, 0.6, 1.2, 8.8, 5, es_sql=False)
    
    # Impacto
    impacto = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.7))
    text_frame = impacto.text_frame
    p = text_frame.paragraphs[0]
    p.text = "IMPACTO: Profesor abre Colab → Run All → Funciona en 20 segundos (sin configuración manual)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_EXITO
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================
    # PERSONA 4 - SLIDE 14: CIERRE TÉCNICO (1 min)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    text_frame = titulo.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Mapeo Teoría → Práctica"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    
    # Columna Programación I
    prog1 = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.3), Inches(3.5))
    text_frame = prog1.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "PROGRAMACIÓN I:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUBTITULO
    
    items_prog1 = [
        "Estructuras de datos → DataFrames",
        "Funciones y modularidad → .copy()",
        "Manejo de archivos → pd.read_csv()",
        "Librerías externas → import",
        "Visualización → Matplotlib",
        "Condicionales → if 'google.colab'",
        "Métodos de strings → .strftime()"
    ]
    
    for item in items_prog1:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.level = 1
    
    # Columna Base de Datos II
    bd2 = slide.shapes.add_textbox(Inches(5.1), Inches(1.3), Inches(4.3), Inches(3.5))
    text_frame = bd2.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "BASE DE DATOS II:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUBTITULO
    
    items_bd2 = [
        "Modelo relacional → SQLite",
        "JOINs → pd.merge() LEFT JOIN",
        "Normalización vs Desnorm. → 1 tabla",
        "SQL: SELECT, GROUP BY, ORDER BY",
        "Funciones agregación → SUM, COUNT",
        "Data Warehousing → OLAP",
        "ACID → sqlite3.connect()"
    ]
    
    for item in items_bd2:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.level = 1
    
    # Pipeline completo
    pipeline = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(0.7))
    text_frame = pipeline.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "CSV → Pandas → Transform → SQLite → SQL → Matplotlib → PNG"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "Pipeline ETL completo, reproducible, escalable"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    
    # Cierre
    cierre = slide.shapes.add_textbox(Inches(2.5), Inches(6.3), Inches(5), Inches(0.8))
    text_frame = cierre.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "¿PREGUNTAS? 🤔"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO
    p.alignment = PP_ALIGN.CENTER
    
    # Guardar
    prs.save('PRESENTACION_ETL_TECNICA.pptx')
    print("✅ Presentación TÉCNICA generada: PRESENTACION_ETL_TECNICA.pptx")
    print("📊 14 diapositivas enfocadas en PROCEDIMIENTOS y CÓDIGO")
    print("⏱️ Distribución: 4 personas × ~3.5 min c/u")
    print("🎯 Coincide con GUIA_PRESENTACION_15MIN.md")

if __name__ == "__main__":
    crear_presentacion()
